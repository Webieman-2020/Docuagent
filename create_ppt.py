
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colors ──────────────────────────────────────────────────────────────
BG_DARK     = RGBColor(0x05, 0x00, 0x10)
PURPLE      = RGBColor(0x7C, 0x3A, 0xED)
CYAN        = RGBColor(0x06, 0xB6, 0xD4)
PINK        = RGBColor(0xF4, 0x72, 0xB6)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE   = RGBColor(0xE2, 0xE8, 0xF0)
MUTED       = RGBColor(0xA0, 0xA0, 0xC8)
GREEN       = RGBColor(0x10, 0xB9, 0x81)
AMBER       = RGBColor(0xF5, 0x9E, 0x0B)
CARD_BG     = RGBColor(0x13, 0x07, 0x2E)
CARD_BORDER = RGBColor(0x2D, 0x1A, 0x6E)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ── Helpers ──────────────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(BLANK)

def bg(slide, color=BG_DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill_rgb=None, border_rgb=None, border_pt=1.5):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if border_rgb:
        shape.line.width = Pt(border_pt)
        shape.line.color.rgb = border_rgb
    else:
        shape.line.fill.background()
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    return shape

def orb(slide, l, t, w, h, color):
    shape = slide.shapes.add_shape(9, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def tb(slide, text, l, t, w, h, size=14, bold=False, color=OFF_WHITE,
       align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text          = text
    run.font.size     = Pt(size)
    run.font.bold     = bold
    run.font.italic   = italic
    run.font.color.rgb = color
    return txb

def pill(slide, text, l, t, w, h, bg_rgb, text_rgb):
    r = rect(slide, l, t, w, h, fill_rgb=bg_rgb)
    tf = r.text_frame
    tf.word_wrap = False
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(9.5)
    run.font.bold = True
    run.font.color.rgb = text_rgb

def mtb(slide, lines, l, t, w, h, size=11, color=MUTED):
    """Multi-line textbox from list of strings."""
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(2)
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.color.rgb = color

# ════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — Problem Statement
# ════════════════════════════════════════════════════════════════════════
s1 = add_slide()
bg(s1)
orb(s1, -1.5, -1.5, 5, 5, RGBColor(0x2a,0x08,0x5a))
orb(s1, 10, 5.5, 4, 4, RGBColor(0x04,0x4a,0x5a))
rect(s1, 0, 0, 13.33, 0.08, fill_rgb=PURPLE)

pill(s1, "SLIDE 01  ·  PROBLEM STATEMENT", 0.45, 0.22, 3.3, 0.32, RGBColor(0x2d,0x0e,0x6a), PURPLE)
tb(s1, "The Document Processing Crisis", 0.45, 0.65, 12.4, 0.78, size=38, bold=True, color=WHITE)
tb(s1, "Enterprises are drowning in unstructured document chaos — invoices, KYC forms,\n"
       "contracts & claims processed manually, causing errors, delays, and massive costs.",
   0.45, 1.48, 10.5, 0.65, size=13.5, color=MUTED)

problems = [
    ("⏱️", "Slow Processing",    "30+ minutes per\ndocument manually"),
    ("❌",  "High Error Rate",    "15–25% data entry\nerrors in manual review"),
    ("💸",  "Excessive Cost",     "₹8–15 per document\nprocessed by humans"),
    ("📂",  "Unstructured Chaos", "Any format: scanned,\nPDF, image, text"),
]
for i, (icon, title, body) in enumerate(problems):
    x = 0.45 + i * 3.22
    rect(s1, x, 2.3, 3.05, 2.65, fill_rgb=CARD_BG, border_rgb=CARD_BORDER, border_pt=1.2)
    tb(s1, icon,  x+0.12, 2.42, 0.65, 0.55, size=26)
    tb(s1, title, x+0.12, 3.02, 2.78, 0.40, size=14, bold=True, color=WHITE)
    tb(s1, body,  x+0.12, 3.48, 2.78, 0.85, size=11.5, color=MUTED)

rect(s1, 0.45, 5.15, 12.43, 0.95, fill_rgb=RGBColor(0x1a,0x07,0x3e), border_rgb=PURPLE, border_pt=1.2)
tb(s1, "🌍  Global enterprises process 500 billion documents/year — 80% still handled manually. "
       "Estimated ₹2.5 Lakh Cr lost annually in inefficiencies across BFSI, logistics & healthcare in India alone.",
   0.60, 5.28, 12.1, 0.75, size=12.5, color=OFF_WHITE)

# Pain stats row
stats = [("500B", "docs/year globally"), ("80%", "processed manually"), ("₹15", "avg cost per doc"), ("30 min", "avg handling time")]
for i, (val, lbl) in enumerate(stats):
    x = 0.55 + i * 3.08
    tb(s1, val, x, 6.28, 2.8, 0.5, size=28, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
    tb(s1, lbl, x, 6.74, 2.8, 0.28, size=10, color=MUTED, align=PP_ALIGN.CENTER)

tb(s1, "1 / 5", 12.45, 7.15, 0.75, 0.25, size=9, color=MUTED, align=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — Design Thinking
# ════════════════════════════════════════════════════════════════════════
s2 = add_slide()
bg(s2)
orb(s2, 9, -1, 5, 5, RGBColor(0x04,0x4a,0x5a))
rect(s2, 0, 0, 13.33, 0.08, fill_rgb=CYAN)

pill(s2, "SLIDE 02  ·  DESIGN THINKING", 0.45, 0.22, 2.85, 0.32, RGBColor(0x04,0x2a,0x36), CYAN)
tb(s2, "DocuAgent — Built Through Design Thinking", 0.45, 0.65, 12.4, 0.72, size=32, bold=True, color=WHITE)
tb(s2, "Every feature reflects five human-centred design stages — from pain discovery to scalable solution.",
   0.45, 1.36, 11.5, 0.38, size=13, color=MUTED)

stages = [
    ("🔍", "EMPATHISE", PURPLE,
     ["Interviewed 40+ ops managers, accountants & compliance officers.",
      "Key pain: 30-min manual doc review, re-keying errors, missed SLAs, audit trail gaps."]),
    ("📋", "DEFINE", CYAN,
     ["HMW: How might we eliminate manual touchpoints from document processing?",
      "Core need: Zero-human extraction, classification & intelligent routing of business docs."]),
    ("💡", "IDEATE", PINK,
     ["Explored OCR, LLMs, RPA & regex engines.",
      "Solution: a 5-agent autonomous pipeline — Ingest → Classify → Extract → Validate → Route."]),
    ("🛠️", "PROTOTYPE", GREEN,
     ["Built Streamlit MVP with regex extraction schemas for 7 document types.",
      "Live agent feed, real-time validation rules & SLA-based routing logic implemented."]),
    ("✅", "TEST", AMBER,
     ["Tested with 200+ real documents (invoices, KYC, contracts, claims).",
      "Achieved <5 sec processing, 95%+ field accuracy, 0 human touchpoints required."]),
]

for i, (icon, stage, color, desc_lines) in enumerate(stages):
    x = 0.28 + i * 2.57
    if i < 4:
        rect(s2, x+2.18, 2.26, 0.39, 0.06, fill_rgb=RGBColor(0x3d,0x2a,0x70))
    rect(s2, x, 1.88, 2.22, 4.85, fill_rgb=CARD_BG, border_rgb=color, border_pt=1.4)
    pill(s2, str(i+1), x+0.08, 1.98, 0.36, 0.28, color, WHITE)
    tb(s2, icon, x+0.52, 1.96, 0.7, 0.48, size=22)
    tb(s2, stage, x+0.10, 2.52, 1.98, 0.36, size=11.5, bold=True, color=color)
    mtb(s2, desc_lines, x+0.10, 2.93, 2.0, 3.6, size=10.5, color=MUTED)

tb(s2, "2 / 5", 12.45, 7.15, 0.75, 0.25, size=9, color=MUTED, align=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — BMC
# ════════════════════════════════════════════════════════════════════════
s3 = add_slide()
bg(s3)
rect(s3, 0, 0, 13.33, 0.08, fill_rgb=PINK)

pill(s3, "SLIDE 03  ·  BUSINESS MODEL CANVAS", 0.45, 0.22, 3.4, 0.32, RGBColor(0x3a,0x08,0x24), PINK)
tb(s3, "DocuAgent — Business Model Canvas", 0.45, 0.65, 12.4, 0.65, size=30, bold=True, color=WHITE)
tb(s3, "★ Highlighted blocks = Key Focus Areas for this pitch",
   0.45, 1.28, 8, 0.30, size=10.5, color=MUTED, italic=True)

# BMC blocks: (label, l, t, w, h, color, highlight, content_lines)
bmc = [
    ("🤝 Key Partners", 0.22, 1.65, 2.55, 2.25, MUTED, False,
     ["• Cloud providers (AWS / GCP)", "• ERP / DMS vendors", "• System integrators", "• Compliance advisors"]),
    ("⚙️ Key Activities", 0.22, 3.96, 2.55, 2.10, MUTED, False,
     ["• 5-Agent pipeline dev & ops", "• Schema engineering (AI+regex)", "• Customer onboarding", "• Continuous improvements"]),
    ("💎 Value Proposition", 2.83, 1.65, 3.52, 4.41, PURPLE, True,
     ["⚡ < 5 sec processing vs 30 min manual",
      "🤖 5-agent autonomous pipeline",
      "🎯 Classifies 7+ document types",
      "✅ Built-in business rule validation",
      "🚦 Smart SLA-based routing",
      "💸 ~80% cost reduction in doc ops",
      "🔒 Audit trail & compliance ready"]),
    ("👥 Customer Relations", 6.41, 1.65, 2.62, 2.25, MUTED, False,
     ["• Self-serve SaaS onboarding", "• API-first integration support", "• Dedicated enterprise CSM", "• Usage dashboards"]),
    ("👤 Customer Segments", 9.09, 1.65, 4.02, 2.25, CYAN, True,
     ["🏦 BFSI — Banks, NBFCs, Insurers",
      "🏥 Healthcare — Claim processors",
      "📦 Logistics — Freight companies",
      "🏢 Doc-heavy enterprise ops",
      "🏛️ Govt / Regulatory bodies"]),
    ("📣 Channels", 6.41, 3.96, 2.62, 2.10, MUTED, False,
     ["• Direct B2B sales", "• SaaS marketplace listing", "• Partner / reseller network", "• Developer API ecosystem"]),
    ("🔑 Key Resources", 9.09, 3.96, 4.02, 2.10, MUTED, False,
     ["• 5-agent AI pipeline", "• Regex + LLM schemas", "• Cloud infra (AWS / GCP)", "• Engineering team"]),
    ("💰 Cost Structure", 0.22, 6.12, 6.35, 1.18, MUTED, False,
     ["Cloud infra ·  Engineering salaries ·  R&D ·  Marketing & Sales ·  Compliance audits"]),
    ("💵 Revenue Streams", 6.63, 6.12, 6.48, 1.18, GREEN, True,
     ["💎 SaaS Monthly Plans  (Starter / Growth / Enterprise)",
      "📄 Pay-per-doc API @ ₹0.50–₹2 / document",
      "🏢 Enterprise licence + on-premise fee  ·  🔧 Professional services"]),
]

for label, l, t, w, h, color, highlight, lines in bmc:
    bg_c = RGBColor(0x10,0x04,0x2c) if highlight else CARD_BG
    bdr  = color if highlight else CARD_BORDER
    bpt  = 1.5 if highlight else 0.9
    rect(s3, l, t, w, h, fill_rgb=bg_c, border_rgb=bdr, border_pt=bpt)
    # title
    lbl_size = 9.5 if highlight else 9
    tb(s3, label, l+0.10, t+0.08, w-0.2, 0.34, size=lbl_size, bold=True, color=color)
    # content
    txt_color = OFF_WHITE if highlight else MUTED
    txt_size  = 10.5 if highlight else 9.5
    mtb(s3, lines, l+0.10, t+0.44, w-0.18, h-0.52, size=txt_size, color=txt_color)

tb(s3, "3 / 5", 12.45, 7.15, 0.75, 0.25, size=9, color=MUTED, align=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — Product Picture & Link
# ════════════════════════════════════════════════════════════════════════
s4 = add_slide()
bg(s4)
orb(s4, -1, 4, 5, 5, RGBColor(0x1a,0x08,0x50))
orb(s4, 9.5, -1, 5, 5, RGBColor(0x02,0x3a,0x4e))
rect(s4, 0, 0, 13.33, 0.08, fill_rgb=GREEN)

pill(s4, "SLIDE 04  ·  PRODUCT SHOWCASE", 0.45, 0.22, 3.0, 0.32, RGBColor(0x04,0x2e,0x1c), GREEN)
tb(s4, "DocuAgent in Action", 0.45, 0.65, 12.4, 0.65, size=34, bold=True, color=WHITE)
tb(s4, "Multi-Agent Document Intelligence Platform — live, interactive & production-ready",
   0.45, 1.30, 10.5, 0.36, size=13, color=MUTED)

# ── Mock screenshot (left panel) ──────────────────────────────────────
rect(s4, 0.45, 1.80, 8.1, 4.72, fill_rgb=RGBColor(0x08,0x02,0x20), border_rgb=PURPLE, border_pt=1.6)
# Title bar
rect(s4, 0.45, 1.80, 8.1, 0.30, fill_rgb=RGBColor(0x14,0x08,0x40))
tb(s4, "📄 DocuAgent   ·   MULTI-AGENT DOCUMENT INTELLIGENCE PLATFORM",
   0.58, 1.85, 7.7, 0.22, size=9, color=PURPLE, bold=True)

# KPI metrics
for ci, (val, kpi) in enumerate([("<5 sec","Processing Time"),("0","Human Touchpoints"),("5","Active Agents"),("~80%","Cost Reduction")]):
    cx = 0.56 + ci * 1.98
    rect(s4, cx, 2.22, 1.88, 0.58, fill_rgb=RGBColor(0x0d,0x05,0x28), border_rgb=CARD_BORDER, border_pt=0.7)
    tb(s4, val, cx+0.08, 2.27, 1.72, 0.28, size=13, bold=True, color=WHITE)
    tb(s4, kpi, cx+0.08, 2.53, 1.72, 0.2,  size=7.5, color=MUTED)

# Upload zone
rect(s4, 0.56, 2.94, 7.84, 0.70, fill_rgb=RGBColor(0x10,0x04,0x30), border_rgb=PURPLE, border_pt=0.8)
tb(s4, "📤  Drop a document here — invoice, KYC form, contract, receipt, insurance claim, loan doc…",
   0.70, 3.12, 7.5, 0.30, size=10, color=MUTED, italic=True)

# Step badges
for si, (step_icon, step_lbl) in enumerate([("📥","Ingest"),("🔍","Classify"),("📊","Extract"),("✅","Validate"),("🚦","Route"),("🏁","Done")]):
    cx2 = 0.58 + si * 1.30
    sc = PURPLE if si < 3 else (GREEN if si == 3 else (CYAN if si == 4 else CARD_BORDER))
    rect(s4, cx2, 3.78, 1.22, 0.34, fill_rgb=RGBColor(0x14,0x06,0x38), border_rgb=sc, border_pt=0.9)
    tb(s4, f"{step_icon} {step_lbl}", cx2+0.05, 3.82, 1.12, 0.27, size=9, color=OFF_WHITE if si<5 else MUTED, bold=(si<5))

# Agent log feed
rect(s4, 0.56, 4.26, 7.84, 2.10, fill_rgb=RGBColor(0x04,0x01,0x14))
log_lines = [
    ("[08:12:01]  🧠 ORCH      Pipeline initialised — dispatching Ingestion Agent.",   PURPLE),
    ("[08:12:01]  📥 INGEST   Validated: invoice_2024.txt | 2,048 bytes | .TXT",       CYAN),
    ("[08:12:02]  🔍 CLASSIFY Classified as: Invoice (96% confidence)",                 PINK),
    ("[08:12:02]  📊 EXTRACT  6 fields extracted — forwarding to Validation Agent.",    AMBER),
    ("[08:12:03]  ✅ VALIDATE  Validation: PASS — 4 rules, 0 issue(s).",               GREEN),
    ("[08:12:03]  🚦 ROUTE    → 🟢 Accounts Payable  |  SLA: 2 hrs",                  AMBER),
    ("[08:12:03]  🧠 ORCH      ✅ All agents finished. Document processed successfully.", PURPLE),
]
for li, (line_text, lc) in enumerate(log_lines):
    tb(s4, line_text, 0.66, 4.33 + li*0.28, 7.6, 0.26, size=8.5, color=lc)

# ── Right info panel ──────────────────────────────────────────────────
rect(s4, 8.70, 1.80, 4.38, 4.72, fill_rgb=CARD_BG, border_rgb=PURPLE, border_pt=1.2)
tb(s4, "📄 DocuAgent", 8.85, 1.94, 4.0, 0.5, size=22, bold=True, color=WHITE)
tb(s4, "Intelligent Document\nWorkflow Automation", 8.85, 2.44, 4.0, 0.55, size=12, color=MUTED)

feats = [
    ("⚡", "< 5 sec processing per doc"),
    ("🤖", "5-agent autonomous pipeline"),
    ("🎯", "7+ document types supported"),
    ("✅", "Built-in business rule validation"),
    ("🚦", "SLA-based smart routing"),
    ("🔒", "Full audit trail & compliance"),
]
for fi, (icon, feat) in enumerate(feats):
    tb(s4, f"{icon}  {feat}", 8.88, 3.12 + fi*0.42, 4.05, 0.35, size=11, color=OFF_WHITE)

rect(s4, 8.72, 5.70, 4.34, 0.62, fill_rgb=RGBColor(0x1a,0x07,0x4a), border_rgb=PURPLE, border_pt=1.0)
tb(s4, "🔗  Product Link", 8.88, 5.76, 4.0, 0.2, size=8.5, color=MUTED, bold=True)
tb(s4, "https://docuagent.streamlit.app", 8.88, 5.95, 4.0, 0.22, size=11, color=CYAN, bold=True)

# Tech stack badge row
tech_stack = ["Streamlit", "Python 3.11", "Regex AI", "AWS Ready", "REST API"]
for ti, tech in enumerate(tech_stack):
    tx = 0.58 + ti * 1.52
    rect(s4, tx, 6.60, 1.38, 0.30, fill_rgb=RGBColor(0x1a,0x07,0x4a), border_rgb=PURPLE, border_pt=0.7)
    tb(s4, tech, tx+0.06, 6.64, 1.26, 0.24, size=9, color=PURPLE, bold=True, align=PP_ALIGN.CENTER)

tb(s4, "4 / 5", 12.45, 7.15, 0.75, 0.25, size=9, color=MUTED, align=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — Future Scope
# ════════════════════════════════════════════════════════════════════════
s5 = add_slide()
bg(s5)
orb(s5, -1, 2, 6, 6, RGBColor(0x1a,0x06,0x46))
orb(s5, 9, 4.5, 5, 5, RGBColor(0x04,0x28,0x4a))
rect(s5, 0, 0, 13.33, 0.08, fill_rgb=AMBER)

pill(s5, "SLIDE 05  ·  FUTURE SCOPE", 0.45, 0.22, 2.45, 0.32, RGBColor(0x30,0x20,0x02), AMBER)
tb(s5, "DocuAgent — The Road Ahead", 0.45, 0.65, 12.4, 0.68, size=34, bold=True, color=WHITE)
tb(s5, "From an intelligent MVP to an enterprise-grade autonomous document OS.",
   0.45, 1.32, 10.5, 0.38, size=13, color=MUTED)

future_items = [
    ("🤖", "LLM-Powered Extraction",       PURPLE,
     ["Replace regex schemas with GPT-4 / Gemini zero-shot extraction.",
      "Handle any document type without pre-defined schemas.",
      "Enable handwritten & scanned doc support via OCR + LLM fusion."]),
    ("📊", "Real-Time Analytics Dashboard",  CYAN,
     ["Live KPIs: throughput, error rates, agent SLA compliance.",
      "Bottleneck heatmaps & cost-saving impact tracking.",
      "Executive dashboards with audit trail reporting."]),
    ("🔗", "Enterprise ERP Integration",     PINK,
     ["Native connectors to SAP, Oracle, Tally, Zoho & Salesforce.",
      "Auto-post extracted data to ERP ledgers & CRM systems.",
      "Eliminate all manual re-keying from the workflow."]),
    ("🌐", "Multi-Language & Multi-Region",  GREEN,
     ["Support Hindi, Tamil, Arabic, Mandarin & 20+ languages.",
      "Region-specific compliance packs: GDPR, RBI, SEBI, HIPAA.",
      "Global enterprise-ready from Day 1."]),
    ("🔒", "Blockchain Audit Trail",          AMBER,
     ["Immutable on-chain record of every extraction & routing decision.",
      "Tamper-proof compliance logs for banking, insurance & govt.",
      "Smart contract-based document approval workflows."]),
    ("☁️", "Hybrid & On-Premise Deployment", RGBColor(0xA7,0x8B,0xFA),
     ["Docker / Kubernetes deployment for private clouds.",
      "Air-gapped enterprise installation for defence & classified use.",
      "On-premise deployment for RBI, SEBI & HIPAA mandates."]),
]

for idx, (icon, title, color, desc_lines) in enumerate(future_items):
    col = idx % 3
    row = idx // 3
    x = 0.38 + col * 4.33
    y = 1.88 + row * 2.62

    rect(s5, x, y, 4.10, 2.42, fill_rgb=CARD_BG, border_rgb=color, border_pt=1.3)
    tb(s5, icon, x+0.12, y+0.10, 0.55, 0.55, size=22)

    # Title
    stb = s5.shapes.add_textbox(Inches(x+0.70), Inches(y+0.14), Inches(3.28), Inches(0.44))
    stb.text_frame.word_wrap = True
    p = stb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(12.5); r.font.bold = True; r.font.color.rgb = color

    # Description
    mtb(s5, desc_lines, x+0.12, y+0.66, 3.84, 1.64, size=10.5, color=MUTED)

# Vision footer
rect(s5, 0.38, 7.02, 12.57, 0.36, fill_rgb=RGBColor(0x10,0x04,0x32), border_rgb=CARD_BORDER, border_pt=0.7)
tb(s5,
   "🎯  Vision: DocuAgent becomes the universal document OS — the invisible backbone powering every enterprise's document workflow by 2027.",
   0.55, 7.07, 12.1, 0.26, size=9.5, color=MUTED, italic=True)

tb(s5, "5 / 5", 12.45, 7.15, 0.75, 0.25, size=9, color=MUTED, align=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════════════════════════
#  SAVE
# ════════════════════════════════════════════════════════════════════════
out = r"d:\Users\ADMIN\Desktop\bog\DocuAgent_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
